import os
from datetime import datetime
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PPTXReporter:
    """用 python-pptx 制作 PowerPoint 报告，替代 HWC 的 report 功能。"""

    def __init__(self):
        self._prs: Optional[Presentation] = None
        self._slides: List[dict] = []  # 收集的幻灯片数据

    # ── 初始化 ──

    def create(self):
        """初始化一个新的空白 Presentation（对应原 create_report）。"""
        self._prs = Presentation()
        # 设置为 16:9
        self._prs.slide_width = Inches(13.333)
        self._prs.slide_height = Inches(7.5)
        self._slides.clear()

    # ── 添加幻灯片数据 ──

    def add_image_slide(self, label: str, image_path: str):
        """记录一张图片幻灯片（对应原 add slide "One Image with Caption"）。"""
        self._slides.append({
            'type': 'image_with_caption',
            'label': label,
            'image_path': image_path,
        })

    def add_image_only_slide(self, label: str, image_path: str):
        """记录一张纯图片幻灯片（对应原 add slide "One Image only"）。"""
        self._slides.append({
            'type': 'image_only',
            'label': label,
            'image_path': image_path,
        })

    # ── 构建 & 导出 ──

    def build(self):
        """根据收集的 _slides 数据构建所有 PowerPoint 页面。"""
        if self._prs is None:
            self.create()

        # 封面页
        self._add_title_slide()

        # 内容页
        for slide_data in self._slides:
            stype = slide_data['type']
            if stype == 'image_with_caption':
                self._add_image_caption_slide(slide_data['label'], slide_data['image_path'])
            elif stype == 'image_only':
                self._add_image_only_slide(slide_data['label'], slide_data['image_path'])

    def export(self, output_path: str) -> str:
        """构建并导出 PPTX 文件。返回输出路径。"""
        self.build()
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self._prs.save(output_path)
        return output_path

    # ── 内部方法 ──

    def _add_title_slide(self):
        """添加封面幻灯片。"""
        slide_layout = self._prs.slide_layouts[6]  # Blank
        slide = self._prs.slides.add_slide(slide_layout)
        slide_w = self._prs.slide_width
        slide_h = self._prs.slide_height

        # 标题
        left = Inches(1)
        top = Inches(2.5)
        width = slide_w - Inches(2)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "HyperView Post-Processing Report"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 56, 100)
        p.alignment = PP_ALIGN.CENTER

        # 时间
        top2 = Inches(3.8)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(0.6))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = datetime.now().strftime("%Y-%m-%d  %H:%M:%S")
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(120, 120, 120)
        p2.alignment = PP_ALIGN.CENTER

    def _add_image_caption_slide(self, label: str, image_path: str):
        """One Image with Caption: 标题 + 大图。"""
        slide_layout = self._prs.slide_layouts[6]  # Blank
        slide = self._prs.slides.add_slide(slide_layout)
        slide_w = self._prs.slide_width
        slide_h = self._prs.slide_height

        # 标题栏
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), slide_w - Inches(1), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 56, 100)

        # 图片居中放置
        if os.path.exists(image_path):
            img_max_w = slide_w - Inches(1)
            img_max_h = slide_h - Inches(1.2)
            img_left, img_top, img_w, img_h = self._fit_image(
                image_path, Inches(0.5), Inches(1.0), img_max_w, img_max_h
            )
            slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)

    def _add_image_only_slide(self, label: str, image_path: str):
        """One Image only: 图片铺满页面，标题在底部。"""
        slide_layout = self._prs.slide_layouts[6]  # Blank
        slide = self._prs.slides.add_slide(slide_layout)
        slide_w = self._prs.slide_width
        slide_h = self._prs.slide_height

        if os.path.exists(image_path):
            img_max_w = slide_w - Inches(0.4)
            img_max_h = slide_h - Inches(0.8)
            img_left, img_top, img_w, img_h = self._fit_image(
                image_path, Inches(0.2), Inches(0.2), img_max_w, img_max_h
            )
            slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)

        # 底部小标签
        txBox = slide.shapes.add_textbox(Inches(0.3), slide_h - Inches(0.5), slide_w - Inches(0.6), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.LEFT

    @staticmethod
    def _fit_image(image_path: str, area_left, area_top, area_w, area_h):
        """计算图片在给定区域内保持比例的最大尺寸和居中位置。

        Returns (left, top, width, height) as Emu values.
        """
        from PIL import Image
        with Image.open(image_path) as img:
            img_w_px, img_h_px = img.size

        # 转成 Emu 计算比例
        area_w_emu = int(area_w)
        area_h_emu = int(area_h)

        scale_w = area_w_emu / img_w_px
        scale_h = area_h_emu / img_h_px
        scale = min(scale_w, scale_h)

        final_w = int(img_w_px * scale)
        final_h = int(img_h_px * scale)

        # 居中
        left = int(area_left) + (area_w_emu - final_w) // 2
        top = int(area_top) + (area_h_emu - final_h) // 2

        return left, top, final_w, final_h
