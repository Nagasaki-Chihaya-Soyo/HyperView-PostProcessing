"""合并多个单 slide PPTX 文件为一个完整的 PPTX。"""
import os
import shutil
from copy import deepcopy
from pptx import Presentation


_R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def merge_pptx(pptx_paths: list, output_path: str) -> str:
    """
    将多个 PPTX（每个通常只包含 1 张 slide）合并为一个。

    Parameters
    ----------
    pptx_paths : list[str]
        各临时 PPTX 文件路径。
    output_path : str
        输出路径。

    Returns
    -------
    str  实际写入的文件路径。
    """
    valid = [p for p in pptx_paths if p and os.path.isfile(p)]
    if not valid:
        return ''

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    if len(valid) == 1:
        shutil.copy2(valid[0], output_path)
        return output_path

    prs = Presentation(valid[0])

    for path in valid[1:]:
        src_prs = Presentation(path)
        for src_slide in src_prs.slides:
            _copy_slide(prs, src_slide)

    prs.save(output_path)
    return output_path


def _copy_slide(target_prs, src_slide):
    """将 src_slide 复制到 target_prs 中。"""
    blank = target_prs.slide_layouts[6]  # blank layout
    new_slide = target_prs.slides.add_slide(blank)

    # 复制非 layout 的 relationships（图片等），建立 rId 映射
    rid_map = {}
    for rel in src_slide.part.rels.values():
        if rel.is_external:
            continue
        if 'slideLayout' in rel.reltype:
            continue
        new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rId

    # 复制 shape 元素，更新其中引用的 rId
    for child in src_slide.shapes._spTree:
        tag_local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag_local in ('nvGrpSpPr', 'grpSpPr'):
            continue  # blank layout 已自带这两个节点
        el = deepcopy(child)
        for node in el.iter():
            for attr in (f'{{{_R_NS}}}embed', f'{{{_R_NS}}}link'):
                old_id = node.get(attr)
                if old_id and old_id in rid_map:
                    node.set(attr, rid_map[old_id])
        new_slide.shapes._spTree.append(el)

    return new_slide
